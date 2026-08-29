import io
import re
import streamlit as st
from pypdf import PdfReader
import docx2txt
import pptx
from langchain_groq import ChatGroq
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings

# -----------------------------------------------------------------------------
# 1. PAGE CONFIGURATION & STYLING
# -----------------------------------------------------------------------------
st.set_page_config(
    page_title="AI Tutor Chatbot",
    page_icon="🤖",
    layout="wide"
)

st.markdown(
    """
    <style>
    .main { background-color: #ffffff; color: #1f2937; }
    .content-box {
        background-color: #f8fafc;
        padding: 20px;
        border-radius: 8px;
        border-top: 4px solid #3b82f6;
        border-left: 1px solid #e2e8f0;
        border-right: 1px solid #e2e8f0;
        border-bottom: 1px solid #e2e8f0;
        margin-bottom: 20px;
    }
    .stButton>button { 
        background-color: #3b82f6 !important; 
        color: #ffffff !important; 
        border-radius: 8px; 
        font-weight: bold; 
    }
    .stButton>button:hover { background-color: #2563eb !important; }
    </style>
    """,
    unsafe_allow_html=True,
)

# -----------------------------------------------------------------------------
# 2. SESSION STATE & HELPER FUNCTIONS
# -----------------------------------------------------------------------------
if "hourly_materials" not in st.session_state or not st.session_state.hourly_materials:
    st.warning("⚠️ No course materials available. Faculty must upload content first.")
    st.stop()

if "chat_history" not in st.session_state:
    st.session_state.chat_history = []
    
if "active_rag_context" not in st.session_state:
    st.session_state.active_rag_context = None

@st.cache_resource
def get_embeddings_model():
    return HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")

def extract_text_from_bytes(file_name, file_bytes):
    """Extracts text from binary file data saved in session state."""
    text = ""
    try:
        if file_name.endswith(".pdf"):
            reader = PdfReader(io.BytesIO(file_bytes))
            for page in reader.pages:
                t = page.extract_text()
                if t: text += t + "\n"
        elif file_name.endswith((".docx", ".doc")):
            text += docx2txt.process(io.BytesIO(file_bytes)) + "\n"
        elif file_name.endswith((".pptx", ".ppt")):
            prs = pptx.Presentation(io.BytesIO(file_bytes))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif file_name.endswith(".txt"):
            text += file_bytes.decode("utf-8") + "\n"
    except Exception:
        pass
    return text

def compile_phase_text(phase_data):
    """Combines notes, assignments, MCQs, and extracted file text into one RAG context."""
    if not phase_data:
        return ""
        
    compiled = []
    if phase_data.get("notes"): compiled.append(f"Notes:\n{phase_data['notes']}")
    if phase_data.get("class_assignment"): compiled.append(f"Assignments:\n{phase_data['class_assignment']}")
    if phase_data.get("mcqs"): compiled.append(f"MCQs/Knowledge Check:\n{phase_data['mcqs']}")
    
    for f_list in ["files", "assignment_files"]:
        for file_info in phase_data.get(f_list, []):
            extracted = extract_text_from_bytes(file_info["name"], file_info["bytes"])
            if extracted:
                compiled.append(f"Document ({file_info['name']}):\n{extracted}")
                
    return "\n\n".join(compiled)

def build_vectorstore(text):
    """Chunks and embeds the compiled text."""
    if not text.strip():
        return None
    clean_text = re.sub(r'\n+', '\n', text)
    clean_text = re.sub(r'\s+', ' ', clean_text).strip()
    
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=150)
    chunks = text_splitter.split_text(clean_text)
    
    if not chunks: return None
    return FAISS.from_texts(chunks, get_embeddings_model())

def chat_with_rag(api_key, vectorstore, query, history):
    """Retrieves context and streams LLM response."""
    if not api_key:
        return "⚠️ Please enter your Groq API Key in the sidebar."
        
    context = ""
    if vectorstore:
        docs = vectorstore.similarity_search(query, k=4)
        context = "\n\n".join([d.page_content for d in docs])
    
    if not context:
        context = "No specific uploaded documents found for this question. Rely on general academic knowledge."

    sys_prompt = f"""
    You are an expert, encouraging AI Tutor helping a student. 
    Use the following course materials to answer the student's question accurately. If the answer is not in the text, you may use external academic knowledge but clarify that you are doing so.
    
    COURSE MATERIALS CONTEXT:
    {context}
    """
    
    try:
        llm = ChatGroq(groq_api_key=api_key, model_name="openai/gpt-oss-120b", temperature=0.3)
        # Format history for LangChain
        messages = [("system", sys_prompt)]
        for msg in history[-4:]: # Keep last 4 interactions for context memory
            messages.append((msg["role"], msg["content"]))
        messages.append(("user", query))
        
        response = llm.invoke(messages)
        return response.content.strip()
    except Exception as e:
        return f"⚠️ Chat Failed: {str(e)}"

# -----------------------------------------------------------------------------
# 3. SIDEBAR & NAVIGATION
# -----------------------------------------------------------------------------
st.sidebar.title("⚙️ Configuration")
groq_api_key = st.sidebar.text_input("Groq API Key", value=st.session_state.get("groq_api_key", ""), type="password")

st.title("🤖 Interactive AI Course Tutor")
st.caption("Select exactly what phase of the course you want to study, and I will read all the specific notes, PDFs, and assignments to answer your questions.")

# Context Selection Layout
c1, c2, c3, c4 = st.columns(4)

with c1:
    subjects = list(st.session_state.hourly_materials.keys())
    selected_sub = st.selectbox("1. Subject", options=subjects)

if selected_sub:
    with c2:
        modules = list(st.session_state.hourly_materials[selected_sub].keys())
        selected_mod = st.selectbox("2. Module", options=modules)
    
    if selected_mod:
        with c3:
            hours = list(st.session_state.hourly_materials[selected_sub][selected_mod].keys())
            selected_hour = st.selectbox("3. Hour", options=hours)
            
        with c4:
            phase_map = {"pre": "🌅 Pre-Class", "main": "🎓 Main Lecture", "post": "🌇 Post-Class"}
            selected_phase = st.selectbox("4. Learning Phase", options=list(phase_map.keys()), format_func=lambda x: phase_map[x])

st.markdown("---")

# -----------------------------------------------------------------------------
# 4. RAG ACTIVATION
# -----------------------------------------------------------------------------
# Generate a unique string identifier for the student's current dropdown selection
current_context_id = f"{selected_sub}_{selected_mod}_{selected_hour}_{selected_phase}"

col_a, col_b = st.columns([1, 4])
with col_a:
    if st.button("Load Tutor Database"):
        with st.spinner("Extracting documents and compiling AI memory..."):
            phase_data = st.session_state.hourly_materials[selected_sub][selected_mod][selected_hour].get(selected_phase, {})
            raw_text = compile_phase_text(phase_data)
            
            if not raw_text.strip():
                st.warning("No text or documents found in this phase to train the tutor on.")
                st.session_state.active_rag_store = None
            else:
                st.session_state.active_rag_store = build_vectorstore(raw_text)
                st.session_state.active_rag_context = current_context_id
                st.session_state.chat_history = [] # Clear chat history on new topic load
                st.success("Tutor Ready!")

with col_b:
    if st.session_state.active_rag_context == current_context_id:
        st.markdown(f"✅ **Active Memory:** Currently tuned into **{selected_sub} ➔ {selected_mod} ➔ {selected_hour} ({phase_map[selected_phase]})**")
    else:
        st.info("👈 Click 'Load Tutor Database' to align the AI's memory with your selected topic above.")

st.markdown("<br>", unsafe_allow_html=True)

# -----------------------------------------------------------------------------
# 5. CHAT INTERFACE
# -----------------------------------------------------------------------------
st.markdown('<div class="content-box">', unsafe_allow_html=True)

# Render Chat History
for msg in st.session_state.chat_history:
    with st.chat_message(msg["role"]):
        st.markdown(msg["content"])

# User Input
if prompt := st.chat_input("Ask a question about the materials in this phase..."):
    if st.session_state.active_rag_context != current_context_id:
        st.error("Please click 'Load Tutor Database' first to initialize the memory for this specific topic.")
    else:
        # Append User Msg
        st.session_state.chat_history.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)

        # Generate & Append Assistant Msg
        with st.chat_message("assistant"):
            with st.spinner("Thinking..."):
                response = chat_with_rag(
                    groq_api_key, 
                    st.session_state.get("active_rag_store"), 
                    prompt, 
                    st.session_state.chat_history
                )
                st.markdown(response)
        st.session_state.chat_history.append({"role": "assistant", "content": response})

st.markdown('</div>', unsafe_allow_html=True)
