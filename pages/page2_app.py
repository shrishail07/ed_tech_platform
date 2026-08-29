

# import io
# import re
# import streamlit as st
# from pypdf import PdfReader
# import docx2txt
# import pptx
# from langchain_groq import ChatGroq
# from langchain_text_splitters import RecursiveCharacterTextSplitter
# from langchain_community.vectorstores import FAISS
# from langchain_community.embeddings import HuggingFaceEmbeddings

# # -----------------------------------------------------------------------------
# # 1. PAGE CONFIGURATION & STYLING
# # -----------------------------------------------------------------------------
# st.set_page_config(
#     page_title="Hourly Content Upload & AI Generator",
#     page_icon="🕒",
#     layout="wide"
# )

# st.markdown(
#     """
#     <style>
#     .main { background-color: #ffffff; color: #1f2937; }
#     .stTextArea textarea, .stTextInput input { 
#         background-color: #f9fafb !important; 
#         color: #111827 !important; 
#         border: 1px solid #d1d5db; 
#     }
#     .stButton>button { 
#         background-color: #2563eb !important; 
#         color: #ffffff !important; 
#         border-radius: 8px; 
#         font-weight: bold; 
#     }
#     .stButton>button:hover { background-color: #1d4ed8 !important; }
#     .content-box {
#         background-color: #f8fafc;
#         padding: 20px;
#         border-radius: 8px;
#         border-top: 4px solid #10b981;
#         border-left: 1px solid #e2e8f0;
#         border-right: 1px solid #e2e8f0;
#         border-bottom: 1px solid #e2e8f0;
#         margin-top: 10px;
#     }
#     .ai-box {
#         background-color: #f5f3ff;
#         padding: 15px;
#         border-radius: 8px;
#         border: 1px dashed #8b5cf6;
#         margin-bottom: 15px;
#     }
#     .status-badge {
#         background-color: #ecfdf5;
#         border: 1px solid #10b981;
#         padding: 10px;
#         border-radius: 6px;
#         color: #065f46;
#         margin-bottom: 15px;
#     }
#     </style>
#     """,
#     unsafe_allow_html=True,
# )

# # -----------------------------------------------------------------------------
# # 2. SESSION STATE, EXTRACTION & ADVANCED RAG FUNCTIONS
# # -----------------------------------------------------------------------------
# if "hourly_materials" not in st.session_state:
#     st.session_state.hourly_materials = {}

# if "all_subjects_data" not in st.session_state or not st.session_state.all_subjects_data:
#     st.warning("⚠️ No syllabus data found. Please go back to the Extractor page (Page 1) and process a syllabus first.")
#     st.stop()

# @st.cache_resource
# def get_embeddings_model():
#     """Loads lightweight HuggingFace embeddings for the RAG Vector Store."""
#     return HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")

# def extract_text_from_files(uploaded_files):
#     """Extracts raw text from PDF, DOCX, PPTX, and TXT files."""
#     extracted_text = ""
#     for f in uploaded_files:
#         try:
#             if f.name.endswith(".pdf"):
#                 reader = PdfReader(io.BytesIO(f.getvalue()))
#                 for page in reader.pages:
#                     t = page.extract_text()
#                     if t: extracted_text += t + "\n"
#             elif f.name.endswith((".docx", ".doc")):
#                 extracted_text += docx2txt.process(io.BytesIO(f.getvalue())) + "\n"
#             elif f.name.endswith((".pptx", ".ppt")):
#                 prs = pptx.Presentation(io.BytesIO(f.getvalue()))
#                 for slide in prs.slides:
#                     for shape in slide.shapes:
#                         if hasattr(shape, "text"):
#                             extracted_text += shape.text + "\n"
#             elif f.name.endswith(".txt"):
#                 extracted_text += f.getvalue().decode("utf-8") + "\n"
#         except Exception as e:
#             st.error(f"Error reading file {f.name}: {e}")
#     return extracted_text

# def clean_and_build_rag(text):
#     """Cleans unstructured text, chunks it, and builds a FAISS Vector Database."""
#     if not text or not text.strip():
#         return None
    
#     # 1. Clean Unstructured Data
#     clean_text = re.sub(r'\n+', '\n', text)
#     clean_text = re.sub(r'\s+', ' ', clean_text).strip()
    
#     # 2. Chunking
#     text_splitter = RecursiveCharacterTextSplitter(
#         chunk_size=1000, 
#         chunk_overlap=200,
#         length_function=len
#     )
#     chunks = text_splitter.split_text(clean_text)
    
#     # 3. Build Vector Store
#     if not chunks:
#         return None
#     embeddings = get_embeddings_model()
#     vectorstore = FAISS.from_texts(chunks, embeddings)
#     return vectorstore

# def generate_ai_text(api_key, prompt_text):
#     """Raw Groq LLM Caller"""
#     if not api_key:
#         st.error("Please enter your Groq API Key in the sidebar.")
#         return ""
#     try:
#         llm = ChatGroq(groq_api_key=api_key, model_name="openai/gpt-oss-120b", temperature=0.2)
#         response = llm.invoke(prompt_text)
#         return response.content.strip()
#     except Exception as e:
#         st.error(f"AI Generation Failed: {e}")
#         return ""

# def rag_generate(api_key, vectorstore, query, system_prompt):
#     """Retrieves relevant chunks from FAISS and generates response via LLM."""
#     context = ""
#     if vectorstore:
#         # Retrieve the top 4 most relevant chunks based on the specific query
#         docs = vectorstore.similarity_search(query, k=4)
#         context = "\n\n".join([d.page_content for d in docs])
    
#     if not context:
#         context = "(No specific documents provided. Rely on general academic knowledge for this topic.)"

#     final_prompt = f"""
#     {system_prompt}
    
#     RAG RETRIEVED CONTEXT (Use this strictly to ground your response):
#     \"\"\"
#     {context}
#     \"\"\"
#     """
#     return generate_ai_text(api_key, final_prompt)

# # -----------------------------------------------------------------------------
# # 3. SIDEBAR & NAVIGATION CONTROLS
# # -----------------------------------------------------------------------------
# st.sidebar.title("⚙️ Configuration")
# groq_api_key = st.sidebar.text_input(
#     "Groq API Key", 
#     value=st.session_state.get("groq_api_key", ""), 
#     type="password"
# )
# st.session_state["groq_api_key"] = groq_api_key

# st.title("🕒 Advanced RAG Content & AI Generator")
# st.caption("Upload any highly unstructured document (PDF, PPTX, DOCX). The system will clean it, build a Vector Database, and generate exact materials.")

# col1, col2, col3 = st.columns(3)

# with col1:
#     subject_names = list(st.session_state.all_subjects_data.keys())
#     selected_subject = st.selectbox("1. Select Subject", options=subject_names)

# active_data = st.session_state.all_subjects_data[selected_subject]
# modules = active_data.get("modules", [])

# if not modules:
#     st.error("No modules found for this subject.")
#     st.stop()

# with col2:
#     module_options = {i: f"Module {mod.get('module_num', i+1)}" for i, mod in enumerate(modules)}
#     selected_mod_idx = st.selectbox("2. Select Module", options=list(module_options.keys()), format_func=lambda x: module_options[x])

# active_module = modules[selected_mod_idx]
# module_name = f"Module {active_module.get('module_num', selected_mod_idx+1)}"
# module_concepts = active_module.get("key_concepts", "No concepts provided.")

# with col3:
#     raw_plan = active_module.get("hourly_plan", "")
#     parsed_hours = [hour.strip() for hour in raw_plan.split("\n") if hour.strip()]
#     if not parsed_hours: 
#         parsed_hours = ["Fallback Hour 1"]
#     selected_hour = st.selectbox("3. Select Teaching Hour", options=parsed_hours)

# state_key = f"{selected_subject}_{module_name}_{selected_hour}".replace(" ", "_")

# # Initialize widget keys
# for k in ["main_notes", "pre_notes", "pre_mcqs", "post_notes", "post_mcqs", "youtube_link", "class_assignment"]:
#     widget_key = f"widget_{k}_{state_key}"
#     if widget_key not in st.session_state:
#         st.session_state[widget_key] = ""
        
# if f"{state_key}_main_extracted_doc" not in st.session_state:
#     st.session_state[f"{state_key}_main_extracted_doc"] = ""

# # -----------------------------------------------------------------------------
# # 4. CONTENT AUTHORING & AI TABS
# # -----------------------------------------------------------------------------
# st.markdown("---")
# st.subheader(f"Authoring: {selected_subject} ➔ {module_name} ➔ {selected_hour}")

# with st.container():
#     st.markdown('<div class="content-box">', unsafe_allow_html=True)
    
#     tab_main, tab_pre, tab_post = st.tabs([
#         "🎓 1. Main Lecture Content (Source)", 
#         "🌅 2. Pre-Class Materials (RAG)", 
#         "🌇 3. Post-Class Materials (RAG)"
#     ])
    
#     # =========================================================================
#     # TAB 1: MAIN LECTURE (UPLOAD & RAG PIPELINE TRIGGER)
#     # =========================================================================
#     with tab_main:
#         st.markdown("#### Step 1: Provide Lecture Content")
        
#         main_files = st.file_uploader(
#             "📎 Upload Lecture Documents (PDF, DOCX, PPTX, TXT)", 
#             type=["pdf", "docx", "doc", "pptx", "ppt", "txt"],
#             accept_multiple_files=True, 
#             key=f"widget_main_files_{state_key}"
#         )
        
#         if main_files:
#             st.session_state[f"{state_key}_main_extracted_doc"] = extract_text_from_files(main_files)
        
#         combined_source_text = f"{st.session_state[f'widget_main_notes_{state_key}']}\n\n{st.session_state[f'{state_key}_main_extracted_doc']}".strip()
        
#         # Build RAG Database dynamically when content is provided
#         rag_vectorstore = None
#         if combined_source_text:
#             with st.spinner("Processing documents into Vector Database..."):
#                 rag_vectorstore = clean_and_build_rag(combined_source_text)
#                 st.markdown(
#                     """
#                     <div class="status-badge">
#                         ✅ <strong>RAG Pipeline Active:</strong> Documents cleaned, chunked, and embedded into FAISS Vector Database.
#                     </div>
#                     """, 
#                     unsafe_allow_html=True
#                 )
        
#         st.markdown('<div class="ai-box">', unsafe_allow_html=True)
#         if st.button("✨ Auto-Generate Deep Lecture Notes from Uploads", key=f"btn_gen_main_{state_key}"):
#             with st.spinner("Retrieving RAG context and generating deep notes..."):
#                 query = f"Core concepts, definitions, and comprehensive details regarding {selected_hour}"
#                 sys_prompt = f"""
#                 You are an expert professor teaching '{selected_subject}'.
#                 Write highly detailed lecture notes (approx. 400 words) for the hour topic: '{selected_hour}'.
#                 Include clear definitions and at least one practical example extracted from the provided context.
#                 """
#                 st.session_state[f"widget_main_notes_{state_key}"] = rag_generate(groq_api_key, rag_vectorstore, query, sys_prompt)
#         st.markdown('</div>', unsafe_allow_html=True)
        
#         st.text_area(
#             "📝 Faculty Lecture Notes / Script", 
#             height=250, 
#             placeholder="Type your own notes or use the RAG generator above...",
#             key=f"widget_main_notes_{state_key}"
#         )
        
#         st.markdown("#### Step 2: Additional Resources & Tasks")
#         st.text_input(
#             "🎥 YouTube Video Link",
#             placeholder="e.g., https://www.youtube.com/watch?v=...",
#             key=f"widget_youtube_link_{state_key}"
#         )
        
#         st.text_area(
#             "📋 Class Assignment",
#             placeholder="Describe any in-class activities, homework, or assignments...",
#             height=120,
#             key=f"widget_class_assignment_{state_key}"
#         )

#     # =========================================================================
#     # TAB 2: PRE-CLASS (RAG NOTES + PREREQUISITE MCQS)
#     # =========================================================================
#     with tab_pre:
#         st.markdown('<div class="ai-box">', unsafe_allow_html=True)
#         st.markdown("#### Pre-Class Reading Instructions")
        
#         if st.button("✨ RAG Generate Pre-Class Notes", key=f"btn_gen_pre_{state_key}"):
#             with st.spinner("Retrieving foundational RAG context..."):
#                 query = f"Prerequisites, basics, introductions, and foundational concepts needed for {selected_hour}"
#                 sys_prompt = f"""
#                 Write a concise (120-150 words) Pre-Class Reading Guide for students attending the lecture on "{selected_hour}".
#                 1. State exact prerequisites they must know based on the context.
#                 2. Highlight key terms they should preview.
#                 3. Give 2 guiding questions.
#                 """
#                 st.session_state[f"widget_pre_notes_{state_key}"] = rag_generate(groq_api_key, rag_vectorstore, query, sys_prompt)
        
#         st.text_area("Pre-Class Notes (Editable)", height=160, key=f"widget_pre_notes_{state_key}")
#         pre_files = st.file_uploader("📎 Optional: Supporting Pre-Class File", accept_multiple_files=True, key=f"pre_files_{state_key}")
        
#         st.markdown("---")
#         st.markdown("#### 🧠 Pre-Class Diagnostic MCQs")
#         pc1, pc2, pc3 = st.columns([1, 1, 2])
#         pre_easy = pc1.number_input("Easy Questions", min_value=0, max_value=10, value=2, key=f"pe_{state_key}")
#         pre_hard = pc2.number_input("Hard Questions", min_value=0, max_value=10, value=1, key=f"ph_{state_key}")
        
#         if pc3.button("✨ RAG Generate Pre-Class MCQs", use_container_width=True, key=f"btn_mcq_pre_{state_key}"):
#             with st.spinner("Generating diagnostic MCQs..."):
#                 query = f"Introductory definitions, prerequisites, and basic facts for {selected_hour}"
#                 sys_prompt = f"""
#                 Create {pre_easy} Easy and {pre_hard} Hard Multiple Choice Questions to test introductory understanding of "{selected_hour}".
#                 Format exactly as:
#                 Q[Number]: [Question text]
#                 A) [Option A]
#                 B) [Option B]
#                 C) [Option C]
#                 D) [Option D]
#                 Correct Answer: [Option Letter]
#                 Explanation: [1-sentence explanation]
#                 """
#                 st.session_state[f"widget_pre_mcqs_{state_key}"] = rag_generate(groq_api_key, rag_vectorstore, query, sys_prompt)
                    
#         st.text_area("Pre-Class MCQs (Editable)", height=220, key=f"widget_pre_mcqs_{state_key}")
#         st.markdown('</div>', unsafe_allow_html=True)

#     # =========================================================================
#     # TAB 3: POST-CLASS (RAG SUMMARY + ASSESSMENT MCQS)
#     # =========================================================================
#     with tab_post:
#         st.markdown('<div class="ai-box">', unsafe_allow_html=True)
#         st.markdown("#### Post-Class Summary & Homework")
        
#         if st.button("✨ RAG Generate Post-Class Summary", key=f"btn_gen_post_{state_key}"):
#             with st.spinner("Retrieving advanced RAG context for summary..."):
#                 query = f"Summary, advanced formulas, applications, and conclusions of {selected_hour}"
#                 sys_prompt = f"""
#                 Write a structured Post-Class Takeaway Summary for students who completed "{selected_hour}":
#                 1. Summary of Key Takeaways (3-4 bullet points extracted from context).
#                 2. A Practical Application or Take-Home Problem.
#                 """
#                 st.session_state[f"widget_post_notes_{state_key}"] = rag_generate(groq_api_key, rag_vectorstore, query, sys_prompt)
            
#         st.text_area("Post-Class Summary / Assignments", height=160, key=f"widget_post_notes_{state_key}")
#         post_files = st.file_uploader("📎 Optional: Post-Class File", accept_multiple_files=True, key=f"post_files_{state_key}")
        
#         st.markdown("---")
#         st.markdown("#### 🧠 Post-Class Mastery MCQs")
#         poc1, poc2, poc3 = st.columns([1, 1, 2])
#         post_easy = poc1.number_input("Easy Questions", min_value=0, max_value=10, value=3, key=f"poe_{state_key}")
#         post_hard = poc2.number_input("Hard Questions", min_value=0, max_value=10, value=2, key=f"poh_{state_key}")
        
#         if poc3.button("✨ RAG Generate Post-Class MCQs", use_container_width=True, key=f"btn_mcq_post_{state_key}"):
#             with st.spinner("Generating mastery MCQs..."):
#                 query = f"Complex details, deep understanding, and advanced examples for {selected_hour}"
#                 sys_prompt = f"""
#                 Create {post_easy} Easy and {post_hard} Hard Multiple Choice Questions to test students' mastery of "{selected_hour}".
#                 Format exactly as:
#                 Q[Number]: [Question text]
#                 A) [Option A]
#                 B) [Option B]
#                 C) [Option C]
#                 D) [Option D]
#                 Correct Answer: [Option Letter]
#                 Explanation: [1-sentence explanation]
#                 """
#                 st.session_state[f"widget_post_mcqs_{state_key}"] = rag_generate(groq_api_key, rag_vectorstore, query, sys_prompt)
            
#         st.text_area("Post-Class MCQs (Editable)", height=220, key=f"widget_post_mcqs_{state_key}")
#         st.markdown('</div>', unsafe_allow_html=True)
    
#     st.markdown("<br/>", unsafe_allow_html=True)
    
#     # =========================================================================
#     # COMMIT CHANGES
#     # =========================================================================
#     if st.button("💾 Commit Changes for this Hour", use_container_width=True):
#         if selected_subject not in st.session_state.hourly_materials:
#             st.session_state.hourly_materials[selected_subject] = {}
#         if module_name not in st.session_state.hourly_materials[selected_subject]:
#             st.session_state.hourly_materials[selected_subject][module_name] = {}
            
#         st.session_state.hourly_materials[selected_subject][module_name][selected_hour] = {
#             "main": {
#                 "notes": st.session_state[f"widget_main_notes_{state_key}"],
#                 "youtube_link": st.session_state[f"widget_youtube_link_{state_key}"],
#                 "class_assignment": st.session_state[f"widget_class_assignment_{state_key}"],
#                 "files": [{"name": f.name, "type": f.type, "bytes": f.getvalue()} for f in (main_files or [])]
#             },
#             "pre": {
#                 "notes": st.session_state[f"widget_pre_notes_{state_key}"],
#                 "mcqs": st.session_state[f"widget_pre_mcqs_{state_key}"],
#                 "files": [{"name": f.name, "type": f.type, "bytes": f.getvalue()} for f in (pre_files or [])]
#             },
#             "post": {
#                 "notes": st.session_state[f"widget_post_notes_{state_key}"],
#                 "mcqs": st.session_state[f"widget_post_mcqs_{state_key}"],
#                 "files": [{"name": f.name, "type": f.type, "bytes": f.getvalue()} for f in (post_files or [])]
#             }
#         }
#         st.success(f"✅ All content (RAG Notes, Files, MCQs) saved successfully for '{selected_hour}' in {selected_subject}!")
        
#     st.markdown('</div>', unsafe_allow_html=True)

# # -----------------------------------------------------------------------------
# # 5. REVIEW SAVED CONTENT
# # -----------------------------------------------------------------------------
# st.markdown("---")
# with st.expander("Review Saved Hourly Materials (Session Overview)"):
#     if not st.session_state.hourly_materials:
#         st.info("No materials saved yet.")
#     else:
#         for sub_key, sub_dict in st.session_state.hourly_materials.items():
#             st.markdown(f"### 📘 {sub_key}")
#             for mod_key, hours_dict in sub_dict.items():
#                 st.markdown(f"**{mod_key}**")
#                 for hr_key, content in hours_dict.items():
#                     if "main" in content:
#                         total_files = len(content['main']['files']) + len(content['pre']['files']) + len(content['post']['files'])
#                         pre_mcq_status = "Generated" if content['pre'].get('mcqs') else "Pending"
#                         post_mcq_status = "Generated" if content['post'].get('mcqs') else "Pending"
#                         st.write(f"- **{hr_key}**: {total_files} files | Pre-MCQs: {pre_mcq_status} | Post-MCQs: {post_mcq_status}")


import io
import re
import streamlit as st
from pypdf import PdfReader
import docx2txt
import pptx
from langchain_groq import ChatGroq
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings

# -----------------------------------------------------------------------------
# 1. PAGE CONFIGURATION & STYLING
# -----------------------------------------------------------------------------
st.set_page_config(
    page_title="Hourly Content Upload & AI Generator",
    page_icon="🕒",
    layout="wide"
)

st.markdown(
    """
    <style>
    .main { background-color: #ffffff; color: #1f2937; }
    .stTextArea textarea, .stTextInput input { 
        background-color: #f9fafb !important; 
        color: #111827 !important; 
        border: 1px solid #d1d5db; 
    }
    .stButton>button { 
        background-color: #2563eb !important; 
        color: #ffffff !important; 
        border-radius: 8px; 
        font-weight: bold; 
    }
    .stButton>button:hover { background-color: #1d4ed8 !important; }
    .content-box {
        background-color: #f8fafc;
        padding: 20px;
        border-radius: 8px;
        border-top: 4px solid #10b981;
        border-left: 1px solid #e2e8f0;
        border-right: 1px solid #e2e8f0;
        border-bottom: 1px solid #e2e8f0;
        margin-top: 10px;
    }
    .ai-box {
        background-color: #f5f3ff;
        padding: 15px;
        border-radius: 8px;
        border: 1px dashed #8b5cf6;
        margin-bottom: 15px;
    }
    .status-badge {
        background-color: #ecfdf5;
        border: 1px solid #10b981;
        padding: 10px;
        border-radius: 6px;
        color: #065f46;
        margin-bottom: 15px;
    }
    </style>
    """,
    unsafe_allow_html=True,
)

# -----------------------------------------------------------------------------
# 2. SESSION STATE, EXTRACTION & ADVANCED RAG FUNCTIONS
# -----------------------------------------------------------------------------
if "hourly_materials" not in st.session_state:
    st.session_state.hourly_materials = {}

if "all_subjects_data" not in st.session_state or not st.session_state.all_subjects_data:
    st.warning("⚠️ No syllabus data found. Please go back to the Extractor page (Page 1) and process a syllabus first.")
    st.stop()

@st.cache_resource
def get_embeddings_model():
    """Loads lightweight HuggingFace embeddings for the RAG Vector Store."""
    return HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")

def extract_text_from_files(uploaded_files):
    """Extracts raw text from PDF, DOCX, PPTX, and TXT files."""
    extracted_text = ""
    for f in uploaded_files:
        try:
            if f.name.endswith(".pdf"):
                reader = PdfReader(io.BytesIO(f.getvalue()))
                for page in reader.pages:
                    t = page.extract_text()
                    if t: extracted_text += t + "\n"
            elif f.name.endswith((".docx", ".doc")):
                extracted_text += docx2txt.process(io.BytesIO(f.getvalue())) + "\n"
            elif f.name.endswith((".pptx", ".ppt")):
                prs = pptx.Presentation(io.BytesIO(f.getvalue()))
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            extracted_text += shape.text + "\n"
            elif f.name.endswith(".txt"):
                extracted_text += f.getvalue().decode("utf-8") + "\n"
        except Exception as e:
            st.error(f"Error reading file {f.name}: {e}")
    return extracted_text

def clean_and_build_rag(text):
    """Cleans unstructured text, chunks it, and builds a FAISS Vector Database."""
    if not text or not text.strip():
        return None
    
    clean_text = re.sub(r'\n+', '\n', text)
    clean_text = re.sub(r'\s+', ' ', clean_text).strip()
    
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000, 
        chunk_overlap=200,
        length_function=len
    )
    chunks = text_splitter.split_text(clean_text)
    
    if not chunks:
        return None
    embeddings = get_embeddings_model()
    vectorstore = FAISS.from_texts(chunks, embeddings)
    return vectorstore

def generate_ai_text(api_key, prompt_text):
    """Raw Groq LLM Caller"""
    if not api_key:
        st.error("Please enter your Groq API Key in the sidebar.")
        return ""
    try:
        llm = ChatGroq(groq_api_key=api_key, model_name="openai/gpt-oss-120b", temperature=0.2)
        response = llm.invoke(prompt_text)
        return response.content.strip()
    except Exception as e:
        st.error(f"AI Generation Failed: {e}")
        return ""

def rag_generate(api_key, vectorstore, query, system_prompt):
    """Retrieves relevant chunks from FAISS and generates response via LLM."""
    context = ""
    if vectorstore:
        docs = vectorstore.similarity_search(query, k=4)
        context = "\n\n".join([d.page_content for d in docs])
    
    if not context:
        context = "(No specific documents provided. Rely on general academic knowledge for this topic.)"

    final_prompt = f"""
    {system_prompt}
    
    RAG RETRIEVED CONTEXT (Use this strictly to ground your response):
    \"\"\"
    {context}
    \"\"\"
    """
    return generate_ai_text(api_key, final_prompt)

# -----------------------------------------------------------------------------
# 3. SIDEBAR & NAVIGATION CONTROLS
# -----------------------------------------------------------------------------
st.sidebar.title("⚙️ Configuration")
groq_api_key = st.sidebar.text_input(
    "Groq API Key", 
    value=st.session_state.get("groq_api_key", ""), 
    type="password"
)
st.session_state["groq_api_key"] = groq_api_key

st.title("🕒 Advanced RAG Content & AI Generator")
st.caption("Upload any highly unstructured document (PDF, PPTX, DOCX). The system will clean it, build a Vector Database, and generate exact materials.")

col1, col2, col3 = st.columns(3)

with col1:
    subject_names = list(st.session_state.all_subjects_data.keys())
    selected_subject = st.selectbox("1. Select Subject", options=subject_names)

active_data = st.session_state.all_subjects_data[selected_subject]
modules = active_data.get("modules", [])

if not modules:
    st.error("No modules found for this subject.")
    st.stop()

with col2:
    module_options = {i: f"Module {mod.get('module_num', i+1)}" for i, mod in enumerate(modules)}
    selected_mod_idx = st.selectbox("2. Select Module", options=list(module_options.keys()), format_func=lambda x: module_options[x])

active_module = modules[selected_mod_idx]
module_name = f"Module {active_module.get('module_num', selected_mod_idx+1)}"
module_concepts = active_module.get("key_concepts", "No concepts provided.")

with col3:
    raw_plan = active_module.get("hourly_plan", "")
    parsed_hours = [hour.strip() for hour in raw_plan.split("\n") if hour.strip()]
    if not parsed_hours: 
        parsed_hours = ["Fallback Hour 1"]
    selected_hour = st.selectbox("3. Select Teaching Hour", options=parsed_hours)

state_key = f"{selected_subject}_{module_name}_{selected_hour}".replace(" ", "_")

# Initialize widget keys
for k in ["main_notes", "pre_notes", "pre_mcqs", "post_notes", "post_mcqs", "youtube_link", "google_form_link", "mcq_link", "class_assignment"]:
    widget_key = f"widget_{k}_{state_key}"
    if widget_key not in st.session_state:
        st.session_state[widget_key] = ""
        
if f"{state_key}_main_extracted_doc" not in st.session_state:
    st.session_state[f"{state_key}_main_extracted_doc"] = ""

# -----------------------------------------------------------------------------
# 4. CONTENT AUTHORING & AI TABS
# -----------------------------------------------------------------------------
st.markdown("---")
st.subheader(f"Authoring: {selected_subject} ➔ {module_name} ➔ {selected_hour}")

with st.container():
    st.markdown('<div class="content-box">', unsafe_allow_html=True)
    
    tab_main, tab_pre, tab_post = st.tabs([
        "🎓 1. Main Lecture Content (Source)", 
        "🌅 2. Pre-Class Materials (RAG)", 
        "🌇 3. Post-Class Materials (RAG)"
    ])
    
    # =========================================================================
    # TAB 1: MAIN LECTURE (UPLOAD & RAG PIPELINE TRIGGER)
    # =========================================================================
    with tab_main:
        st.markdown("#### Step 1: Provide Lecture Content")
        
        main_files = st.file_uploader(
            "📎 Upload Lecture Documents (PDF, DOCX, PPTX, TXT)", 
            type=["pdf", "docx", "doc", "pptx", "ppt", "txt"],
            accept_multiple_files=True, 
            key=f"widget_main_files_{state_key}"
        )
        
        if main_files:
            st.session_state[f"{state_key}_main_extracted_doc"] = extract_text_from_files(main_files)
        
        combined_source_text = f"{st.session_state[f'widget_main_notes_{state_key}']}\n\n{st.session_state[f'{state_key}_main_extracted_doc']}".strip()
        
        # Build RAG Database dynamically when content is provided
        rag_vectorstore = None
        if combined_source_text:
            with st.spinner("Processing documents into Vector Database..."):
                rag_vectorstore = clean_and_build_rag(combined_source_text)
                st.markdown(
                    """
                    <div class="status-badge">
                        ✅ <strong>RAG Pipeline Active:</strong> Documents cleaned, chunked, and embedded into FAISS Vector Database.
                    </div>
                    """, 
                    unsafe_allow_html=True
                )
        
        st.markdown('<div class="ai-box">', unsafe_allow_html=True)
        if st.button("✨ Auto-Generate Deep Lecture Notes from Uploads", key=f"btn_gen_main_{state_key}"):
            with st.spinner("Retrieving RAG context and generating deep notes..."):
                query = f"Core concepts, definitions, and comprehensive details regarding {selected_hour}"
                sys_prompt = f"""
                You are an expert professor teaching '{selected_subject}'.
                Write highly detailed lecture notes (approx. 400 words) for the hour topic: '{selected_hour}'.
                Include clear definitions and at least one practical example extracted from the provided context.
                """
                st.session_state[f"widget_main_notes_{state_key}"] = rag_generate(groq_api_key, rag_vectorstore, query, sys_prompt)
        st.markdown('</div>', unsafe_allow_html=True)
        
        st.text_area(
            "📝 Faculty Lecture Notes / Script", 
            height=250, 
            placeholder="Type your own notes or use the RAG generator above...",
            key=f"widget_main_notes_{state_key}"
        )
        
        st.markdown("---")
        st.markdown("#### Step 2: Class Assignments & External Links")
        st.caption("Provide external resources, assignment documents, and assessment links for the students.")
        
        c_link1, c_link2, c_link3 = st.columns(3)
        with c_link1:
            st.text_input(
                "🎥 YouTube Video Link",
                placeholder="https://www.youtube.com/watch?v=...",
                key=f"widget_youtube_link_{state_key}"
            )
        with c_link2:
            st.text_input(
                "📝 Google Form / Survey Link",
                placeholder="https://forms.gle/...",
                key=f"widget_google_form_link_{state_key}"
            )
        with c_link3:
            st.text_input(
                "🔗 External MCQ / Quiz Link",
                placeholder="e.g., Quizizz, Kahoot, or Form link",
                key=f"widget_mcq_link_{state_key}"
            )
        
        st.text_area(
            "📋 Class Assignment Instructions",
            placeholder="Describe any in-class activities, homework, or instructions...",
            height=100,
            key=f"widget_class_assignment_{state_key}"
        )
        
        assignment_files = st.file_uploader(
            "📎 Upload Class Assignment Documents (PDF, DOCX, TXT)", 
            type=["pdf", "docx", "doc", "txt", "pptx"],
            accept_multiple_files=True, 
            key=f"widget_assignment_files_{state_key}"
        )

    # =========================================================================
    # TAB 2: PRE-CLASS (RAG NOTES + PREREQUISITE MCQS)
    # =========================================================================
    with tab_pre:
        st.markdown('<div class="ai-box">', unsafe_allow_html=True)
        st.markdown("#### Pre-Class Reading Instructions")
        
        if st.button("✨ RAG Generate Pre-Class Notes", key=f"btn_gen_pre_{state_key}"):
            with st.spinner("Retrieving foundational RAG context..."):
                query = f"Prerequisites, basics, introductions, and foundational concepts needed for {selected_hour}"
                sys_prompt = f"""
                Write a concise (120-150 words) Pre-Class Reading Guide for students attending the lecture on "{selected_hour}".
                1. State exact prerequisites they must know based on the context.
                2. Highlight key terms they should preview.
                3. Give 2 guiding questions.
                """
                st.session_state[f"widget_pre_notes_{state_key}"] = rag_generate(groq_api_key, rag_vectorstore, query, sys_prompt)
        
        st.text_area("Pre-Class Notes (Editable)", height=160, key=f"widget_pre_notes_{state_key}")
        pre_files = st.file_uploader("📎 Optional: Supporting Pre-Class File", accept_multiple_files=True, key=f"pre_files_{state_key}")
        
        st.markdown("---")
        st.markdown("#### 🧠 Pre-Class Diagnostic MCQs")
        pc1, pc2, pc3 = st.columns([1, 1, 2])
        pre_easy = pc1.number_input("Easy Questions", min_value=0, max_value=10, value=2, key=f"pe_{state_key}")
        pre_hard = pc2.number_input("Hard Questions", min_value=0, max_value=10, value=1, key=f"ph_{state_key}")
        
        if pc3.button("✨ RAG Generate Pre-Class MCQs", use_container_width=True, key=f"btn_mcq_pre_{state_key}"):
            with st.spinner("Generating diagnostic MCQs..."):
                query = f"Introductory definitions, prerequisites, and basic facts for {selected_hour}"
                sys_prompt = f"""
                Create {pre_easy} Easy and {pre_hard} Hard Multiple Choice Questions to test introductory understanding of "{selected_hour}".
                Format exactly as:
                Q[Number]: [Question text]
                A) [Option A]
                B) [Option B]
                C) [Option C]
                D) [Option D]
                Correct Answer: [Option Letter]
                Explanation: [1-sentence explanation]
                """
                st.session_state[f"widget_pre_mcqs_{state_key}"] = rag_generate(groq_api_key, rag_vectorstore, query, sys_prompt)
                    
        st.text_area("Pre-Class MCQs (Editable)", height=220, key=f"widget_pre_mcqs_{state_key}")
        st.markdown('</div>', unsafe_allow_html=True)

    # =========================================================================
    # TAB 3: POST-CLASS (RAG SUMMARY + ASSESSMENT MCQS)
    # =========================================================================
    with tab_post:
        st.markdown('<div class="ai-box">', unsafe_allow_html=True)
        st.markdown("#### Post-Class Summary & Homework")
        
        if st.button("✨ RAG Generate Post-Class Summary", key=f"btn_gen_post_{state_key}"):
            with st.spinner("Retrieving advanced RAG context for summary..."):
                query = f"Summary, advanced formulas, applications, and conclusions of {selected_hour}"
                sys_prompt = f"""
                Write a structured Post-Class Takeaway Summary for students who completed "{selected_hour}":
                1. Summary of Key Takeaways (3-4 bullet points extracted from context).
                2. A Practical Application or Take-Home Problem.
                """
                st.session_state[f"widget_post_notes_{state_key}"] = rag_generate(groq_api_key, rag_vectorstore, query, sys_prompt)
            
        st.text_area("Post-Class Summary / Assignments", height=160, key=f"widget_post_notes_{state_key}")
        post_files = st.file_uploader("📎 Optional: Post-Class File", accept_multiple_files=True, key=f"post_files_{state_key}")
        
        st.markdown("---")
        st.markdown("#### 🧠 Post-Class Mastery MCQs")
        poc1, poc2, poc3 = st.columns([1, 1, 2])
        post_easy = poc1.number_input("Easy Questions", min_value=0, max_value=10, value=3, key=f"poe_{state_key}")
        post_hard = poc2.number_input("Hard Questions", min_value=0, max_value=10, value=2, key=f"poh_{state_key}")
        
        if poc3.button("✨ RAG Generate Post-Class MCQs", use_container_width=True, key=f"btn_mcq_post_{state_key}"):
            with st.spinner("Generating mastery MCQs..."):
                query = f"Complex details, deep understanding, and advanced examples for {selected_hour}"
                sys_prompt = f"""
                Create {post_easy} Easy and {post_hard} Hard Multiple Choice Questions to test students' mastery of "{selected_hour}".
                Format exactly as:
                Q[Number]: [Question text]
                A) [Option A]
                B) [Option B]
                C) [Option C]
                D) [Option D]
                Correct Answer: [Option Letter]
                Explanation: [1-sentence explanation]
                """
                st.session_state[f"widget_post_mcqs_{state_key}"] = rag_generate(groq_api_key, rag_vectorstore, query, sys_prompt)
            
        st.text_area("Post-Class MCQs (Editable)", height=220, key=f"widget_post_mcqs_{state_key}")
        st.markdown('</div>', unsafe_allow_html=True)
    
    st.markdown("<br/>", unsafe_allow_html=True)
    
    # =========================================================================
    # COMMIT CHANGES
    # =========================================================================
    if st.button("💾 Commit Changes for this Hour", use_container_width=True):
        if selected_subject not in st.session_state.hourly_materials:
            st.session_state.hourly_materials[selected_subject] = {}
        if module_name not in st.session_state.hourly_materials[selected_subject]:
            st.session_state.hourly_materials[selected_subject][module_name] = {}
            
        st.session_state.hourly_materials[selected_subject][module_name][selected_hour] = {
            "main": {
                "notes": st.session_state[f"widget_main_notes_{state_key}"],
                "youtube_link": st.session_state[f"widget_youtube_link_{state_key}"],
                "google_form_link": st.session_state[f"widget_google_form_link_{state_key}"],
                "mcq_link": st.session_state[f"widget_mcq_link_{state_key}"],
                "class_assignment": st.session_state[f"widget_class_assignment_{state_key}"],
                "assignment_files": [{"name": f.name, "type": f.type, "bytes": f.getvalue()} for f in (assignment_files or [])],
                "files": [{"name": f.name, "type": f.type, "bytes": f.getvalue()} for f in (main_files or [])]
            },
            "pre": {
                "notes": st.session_state[f"widget_pre_notes_{state_key}"],
                "mcqs": st.session_state[f"widget_pre_mcqs_{state_key}"],
                "files": [{"name": f.name, "type": f.type, "bytes": f.getvalue()} for f in (pre_files or [])]
            },
            "post": {
                "notes": st.session_state[f"widget_post_notes_{state_key}"],
                "mcqs": st.session_state[f"widget_post_mcqs_{state_key}"],
                "files": [{"name": f.name, "type": f.type, "bytes": f.getvalue()} for f in (post_files or [])]
            }
        }
        st.success(f"✅ All content (RAG Notes, Files, MCQs, and Links) saved successfully for '{selected_hour}' in {selected_subject}!")
        
    st.markdown('</div>', unsafe_allow_html=True)

# -----------------------------------------------------------------------------
# 5. REVIEW SAVED CONTENT
# -----------------------------------------------------------------------------
st.markdown("---")
with st.expander("Review Saved Hourly Materials (Session Overview)"):
    if not st.session_state.hourly_materials:
        st.info("No materials saved yet.")
    else:
        for sub_key, sub_dict in st.session_state.hourly_materials.items():
            st.markdown(f"### 📘 {sub_key}")
            for mod_key, hours_dict in sub_dict.items():
                st.markdown(f"**{mod_key}**")
                for hr_key, content in hours_dict.items():
                    if "main" in content:
                        total_files = len(content['main']['files']) + len(content['pre']['files']) + len(content['post']['files'])
                        pre_mcq_status = "Generated" if content['pre'].get('mcqs') else "Pending"
                        post_mcq_status = "Generated" if content['post'].get('mcqs') else "Pending"
                        st.write(f"- **{hr_key}**: {total_files} files | Pre-MCQs: {pre_mcq_status} | Post-MCQs: {post_mcq_status}")
